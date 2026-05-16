"""
PPT generation skill — creates PowerPoint presentations for thesis defense,
proposal, and mid-term reports using python-pptx.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import Optional

# Academic color palette
PRIMARY = RGBColor(0x1A, 0x56, 0x8E)       # Deep blue
SECONDARY = RGBColor(0x3A, 0x7C, 0xA5)     # Medium blue
ACCENT = RGBColor(0xF4, 0xF8, 0xFB)        # Light blue-gray
DARK = RGBColor(0x2D, 0x2D, 0x2D)          # Near black text
MUTED = RGBColor(0x7A, 0x8A, 0x9A)         # Gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xD0, 0xD8, 0xE0)

TITLE_FONT_SIZE = Pt(32)
HEADING_FONT_SIZE = Pt(24)
BODY_FONT_SIZE = Pt(16)
SMALL_FONT_SIZE = Pt(12)
FONT_NAME = "Microsoft YaHei"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=BODY_FONT_SIZE,
                  color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_frame(slide, left, top, width, height, items, font_size=BODY_FONT_SIZE,
                      color=DARK, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = font_name
        p.level = 0
        p.space_after = Pt(6)
    return tf


def _add_accent_bar(slide, left, top, width, height, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_pptx(slides_data: list[dict], output_path: str, title: str = "论文PPT",
                subtitle: str = "", author: str = "") -> str:
    """
    Generate a thesis PPTX file.

    slides_data: list of dicts, each with:
        - layout: "title" | "toc" | "content" | "section" | "reference" | "thanks"
        - title: slide title
        - items: list of bullet strings (for content layout)
        - note: optional note text

    Returns the output_path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # blank

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        layout = sd.get("layout", "content")

        if layout == "title":
            _set_slide_bg(slide, WHITE)
            _add_accent_bar(slide, Inches(1), Inches(2.8), Inches(1.5), Pt(6), PRIMARY)
            _add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.2),
                          sd.get("title", title), TITLE_FONT_SIZE, PRIMARY, bold=True)
            _add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                          sd.get("subtitle", subtitle), BODY_FONT_SIZE, MUTED)
            _add_text_box(slide, Inches(1), Inches(5.5), Inches(5), Inches(0.5),
                          author, SMALL_FONT_SIZE, MUTED)

        elif layout == "toc":
            _set_slide_bg(slide, WHITE)
            _add_accent_bar(slide, Inches(1), Inches(1.2), Inches(0.8), Pt(5), PRIMARY)
            _add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                          sd.get("title", "目录"), HEADING_FONT_SIZE, DARK, bold=True)
            items = sd.get("items", [])
            if items:
                _add_bullet_frame(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(4.5),
                                  items, BODY_FONT_SIZE, DARK)

        elif layout == "section":
            _set_slide_bg(slide, PRIMARY)
            _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
                          sd.get("section_num", ""), Pt(18), RGBColor(0xAA, 0xCC, 0xEE))
            _add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.5),
                          sd.get("title", ""), Pt(36), WHITE, bold=True)
            _add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
                          sd.get("subtitle", ""), BODY_FONT_SIZE, RGBColor(0xCC, 0xDD, 0xEE))

        elif layout == "reference":
            _set_slide_bg(slide, WHITE)
            _add_accent_bar(slide, Inches(1), Inches(1.2), Inches(0.8), Pt(5), PRIMARY)
            _add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                          sd.get("title", "参考文献"), HEADING_FONT_SIZE, DARK, bold=True)
            refs = sd.get("items", [])
            if refs:
                _add_bullet_frame(slide, Inches(1), Inches(2.5), Inches(11), Inches(4.5),
                                  refs, SMALL_FONT_SIZE, DARK)

        elif layout == "thanks":
            _set_slide_bg(slide, PRIMARY)
            _add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
                          sd.get("title", "谢谢！"), Pt(44), WHITE, bold=True,
                          alignment=PP_ALIGN.CENTER)
            _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
                          sd.get("subtitle", "恳请各位老师批评指正"), Pt(20),
                          RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

        else:  # content (default)
            _set_slide_bg(slide, WHITE)
            # Top accent line
            _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), PRIMARY)
            # Title
            _add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
                          sd.get("title", ""), HEADING_FONT_SIZE, DARK, bold=True)
            # Thin divider
            _add_accent_bar(slide, Inches(1), Inches(1.35), Inches(3), Pt(2), ACCENT)
            # Content bullets
            items = sd.get("items", [])
            if items:
                _add_bullet_frame(slide, Inches(1), Inches(1.7), Inches(11.3), Inches(5.0),
                                  items, BODY_FONT_SIZE, DARK)
            # Optional note
            note = sd.get("note", "")
            if note:
                _add_text_box(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.5),
                              note, SMALL_FONT_SIZE, MUTED)

    prs.save(output_path)
    return output_path


def build_thesis_ppt(content_blocks: dict, output_path: str,
                     title: str = "论文答辩", author: str = "") -> str:
    """
    High-level helper: builds a complete thesis PPT from structured content blocks.

    content_blocks expects:
        {
            "title": "论文题目",
            "author": "姓名",
            "advisor": "导师",
            "sections": [
                {"title": "研究背景", "items": ["要点1", "要点2", ...]},
                ...
            ],
            "references": ["[1] 作者. 标题...", ...],
        }
    """
    slides = []

    # Title slide
    slides.append({
        "layout": "title",
        "title": content_blocks.get("title", title),
        "subtitle": content_blocks.get("subtitle", ""),
    })

    # TOC
    toc_items = [s["title"] for s in content_blocks.get("sections", [])]
    if toc_items:
        slides.append({"layout": "toc", "title": "目录", "items": toc_items})

    # Section divider + content slides
    for sec in content_blocks.get("sections", []):
        slides.append({
            "layout": "section",
            "title": sec.get("title", ""),
            "subtitle": sec.get("subtitle", ""),
        })
        slides.append({
            "layout": "content",
            "title": sec.get("title", ""),
            "items": sec.get("items", []),
            "note": sec.get("note", ""),
        })

    # References
    refs = content_blocks.get("references", [])
    if refs:
        slides.append({"layout": "reference", "title": "参考文献", "items": refs})

    # Thanks
    slides.append({"layout": "thanks", "title": "谢谢！", "subtitle": "恳请各位老师批评指正"})

    return create_pptx(slides, output_path, title=title, author=author)
