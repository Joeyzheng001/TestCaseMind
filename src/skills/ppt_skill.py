"""
Simple PPTX generation from slide outlines for thesis defense.

Uses python-pptx to create native editable slides from structured outlines.
For advanced SVG-to-PPTX pipeline, use ppt_engine directly.
"""

from __future__ import annotations

import os
from typing import Dict, List, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Layout constants ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Academic color palette
COLOR_PRIMARY = RGBColor(0x1A, 0x3C, 0x6E)  # Deep navy
COLOR_ACCENT = RGBColor(0xC4, 0x2B, 0x2B)  # Academic red
COLOR_BG_DARK = RGBColor(0x1A, 0x3C, 0x6E)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
COLOR_LINE = RGBColor(0x1A, 0x3C, 0x6E)


def create_pptx(
    slides_data: List[Dict[str, Any]],
    output_path: str,
    title: str = "",
    subtitle: str = "",
) -> str:
    """Create a PPTX file from a list of slide definitions.

    Each slide dict has:
        layout: "title" | "toc" | "content" | "thanks" | "section" | "blank"
        title: Slide title
        subtitle: Subtitle (title/thanks/section layouts)
        items: List of bullet point strings (toc/content layouts)
        image: Optional image path (content layout)

    Args:
        slides_data: List of slide definition dicts.
        output_path: Path to write the .pptx file.
        title: Overall presentation title.
        subtitle: Overall presentation subtitle.

    Returns:
        The output_path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # blank

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        slide = prs.slides.add_slide(blank_layout)

        if layout == "title":
            _build_title_slide(slide, slide_data, title, subtitle)
        elif layout == "toc":
            _build_toc_slide(slide, slide_data)
        elif layout == "section":
            _build_section_slide(slide, slide_data)
        elif layout == "thanks":
            _build_thanks_slide(slide, slide_data)
        else:
            _build_content_slide(slide, slide_data)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _build_title_slide(slide, data, title, subtitle):
    """Title/cover slide with full-bleed color background."""
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    # Accent bar at bottom
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.3), SLIDE_WIDTH, Inches(0.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    slide_title = data.get("title", title)
    slide_subtitle = data.get("subtitle", subtitle)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.3), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title or "论文答辩"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if slide_subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = slide_subtitle
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p2.alignment = PP_ALIGN.CENTER


def _build_toc_slide(slide, data):
    """Table of contents / agenda slide."""
    _add_slide_header(slide, data.get("title", "目录"))

    items = data.get("items", [])
    y_start = Inches(1.8)
    for i, item in enumerate(items):
        y = y_start + Inches(i * 0.65)
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Item text
        txBox = slide.shapes.add_textbox(Inches(2.3), y, Inches(9.5), Inches(0.5))
        tf2 = txBox.text_frame
        p = tf2.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT


def _build_content_slide(slide, data):
    """Content slide with title and bullet points."""
    _add_slide_header(slide, data.get("title", ""))

    items = data.get("items", [])
    y_start = Inches(1.7)
    txBox = slide.shapes.add_textbox(Inches(1.5), y_start, Inches(10.3), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(12)
        p.level = 0
        # Bullet
        p.font.name = "Microsoft YaHei"

    # Bottom accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), SLIDE_HEIGHT - Inches(0.6), Inches(0.8), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()


def _build_section_slide(slide, data):
    """Section divider slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    section_title = data.get("title", "")
    section_subtitle = data.get("subtitle", "")

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    if section_subtitle:
        p2 = tf.add_paragraph()
        p2.text = section_subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p2.alignment = PP_ALIGN.CENTER


def _build_thanks_slide(slide, data):
    """Thank you / Q&A slide."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("title", "谢谢！")
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = data.get("subtitle", "恳请各位老师批评指正")
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER


def _add_slide_header(slide, title_text):
    """Add a standard slide header bar with title."""
    # Top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # Accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.1), SLIDE_WIDTH, Inches(0.04)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_ACCENT
    stripe.line.fill.background()

    if title_text:
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.2), Inches(10.3), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
